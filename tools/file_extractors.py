"""Multi-format extractors -> normalized (text_blocks, tables) per "page".

Supported: PDF, images (RapidOCR), Word (.docx), PowerPoint (.pptx),
spreadsheets (.xlsx/.xls/.csv), plain text (.txt/.md).

Each extractor returns a list of "page" dicts:
    {"page": int, "blocks": [str, ...], "tables": [[[cell,...],...], ...],
     "ocr": bool}
so the extraction agent can treat every format uniformly.
"""
from __future__ import annotations

import os

from tools import ocr_tool
from tools.pdf_extractor import extract_text_blocks, page_is_empty
from tools.table_extractor import extract_tables

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp", ".gif"}
DOC_EXTS = {".docx"}
PPT_EXTS = {".pptx"}
SHEET_EXTS = {".xlsx", ".xls", ".csv"}
TEXT_EXTS = {".txt", ".md", ".markdown"}
PDF_EXTS = {".pdf"}

SUPPORTED_EXTS = (PDF_EXTS | IMAGE_EXTS | DOC_EXTS | PPT_EXTS | SHEET_EXTS | TEXT_EXTS)

def extract_any(path: str) -> list[dict]:
    """Dispatch on extension. Returns list of page dicts (see module docstring)."""
    ext = os.path.splitext(path)[1].lower()
    if ext in PDF_EXTS:
        return _pdf(path)
    if ext in IMAGE_EXTS:
        return _image(path)
    if ext in DOC_EXTS:
        return _docx(path)
    if ext in PPT_EXTS:
        return _pptx(path)
    if ext in SHEET_EXTS:
        return _sheet(path)
    if ext in TEXT_EXTS:
        return _text(path)
    try:
        return _text(path)
    except Exception:
        return []

def _pdf(path: str) -> list[dict]:
    text_pages = extract_text_blocks(path)
    tables_by_page = extract_tables(path)
    out = []
    for rec in text_pages:
        blocks = list(rec["blocks"])
        used_ocr = False
        if page_is_empty(rec) and ocr_tool.ocr_available():
            txt = ocr_tool.ocr_page(path, rec["page"])
            if txt:
                blocks = [txt]
                used_ocr = True
        out.append({"page": rec["page"], "blocks": blocks,
                    "tables": tables_by_page.get(rec["page"], []), "ocr": used_ocr})
    return out

def _image(path: str) -> list[dict]:
    if not ocr_tool.ocr_available():
        return [{"page": 1, "blocks": [], "tables": [], "ocr": False,
                 "_warn": "image needs OCR (pip install rapidocr-onnxruntime)"}]
    text = ocr_tool.ocr_image(path)
    return [{"page": 1, "blocks": [text] if text else [], "tables": [], "ocr": bool(text)}]

def _docx(path: str) -> list[dict]:
    import docx
    d = docx.Document(path)
    blocks = [p.text.strip() for p in d.paragraphs if p.text.strip()]
    tables = []
    for t in d.tables:
        rows = [[c.text.strip() for c in r.cells] for r in t.rows]
        if rows:
            tables.append(rows)
    return [{"page": 1, "blocks": blocks, "tables": tables, "ocr": False}]

def _pptx(path: str) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        blocks = [s.text_frame.text.strip() for s in slide.shapes
                  if s.has_text_frame and s.text_frame.text.strip()]
        out.append({"page": i, "blocks": blocks, "tables": [], "ocr": False})
    return out or [{"page": 1, "blocks": [], "tables": [], "ocr": False}]

def _sheet(path: str) -> list[dict]:
    import pandas as pd
    ext = os.path.splitext(path)[1].lower()
    if ext == ".csv":
        frames = {"Sheet1": pd.read_csv(path, dtype=str, keep_default_na=False)}
    else:
        frames = {k: v.fillna("") for k, v in pd.read_excel(path, sheet_name=None, dtype=str).items()}
    out = []
    for i, (name, df) in enumerate(frames.items(), start=1):
        header = [str(c) for c in df.columns]
        rows = [header] + df.astype(str).values.tolist()
        out.append({"page": i, "blocks": [f"Sheet: {name}"], "tables": [rows], "ocr": False})
    return out or [{"page": 1, "blocks": [], "tables": [], "ocr": False}]

def _text(path: str) -> list[dict]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    blocks = [b.strip() for b in content.split("\n\n") if b.strip()]
    return [{"page": 1, "blocks": blocks, "tables": [], "ocr": False}]
